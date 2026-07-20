import asyncio
import json
import logging
import sys

import ftfy
import requests
from azure.identity import DefaultAzureCredential
from langchain_community.document_loaders import (
    AzureAIDocumentIntelligenceLoader,
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    PyPDFLoader,
    TextLoader,
    YoutubeLoader,
)
from langchain_core.documents import Document
from open_webui.env import (
    AIOHTTP_CLIENT_SESSION_SSL,
    CONTENT_EXTRACTION_MAX_CONCURRENCY,
    CONTENT_EXTRACTION_MAX_OUTPUT_CHARS,
    CONTENT_EXTRACTION_REQUEST_TIMEOUT,
    CONTENT_EXTRACTION_TIMEOUT,
    GLOBAL_LOG_LEVEL,
    RAG_IMAGE_VISION_LLM_API_KEY,
    RAG_IMAGE_VISION_LLM_BASE_URL,
    RAG_IMAGE_VISION_LLM_COMBINE_OCR,
    RAG_IMAGE_VISION_LLM_EXTRA_BODY,
    RAG_IMAGE_VISION_LLM_MAX_TOKENS,
    RAG_IMAGE_VISION_LLM_MODEL,
    RAG_IMAGE_VISION_LLM_PROMPT,
    RAG_PDF_FAST_PATH_MIN_CHARS_PER_PAGE,
    RAG_PDF_IMAGE_MIN_COUNT,
    RAG_PDF_IMAGE_MIN_PIXELS,
    RAG_PDF_IMAGE_ROUTE_ENABLED,
    REQUESTS_VERIFY,
)
from open_webui.retrieval.loaders.datalab_marker import DatalabMarkerLoader
from open_webui.retrieval.loaders.external_document import ExternalDocumentLoader
from open_webui.retrieval.loaders.extraction_ab import get_state
from open_webui.retrieval.loaders.mineru import MinerULoader
from open_webui.retrieval.loaders.mistral import MistralLoader
from open_webui.retrieval.loaders.paddleocr_vl import PaddleOCRVLLoader

logging.basicConfig(stream=sys.stdout, level=GLOBAL_LOG_LEVEL)
log = logging.getLogger(__name__)

# Bounds concurrent document extractions per process (see env.py
# CONTENT_EXTRACTION_MAX_CONCURRENCY). Acquired in Loader.aload so an upload burst
# can't exhaust the default asyncio.to_thread worker pool.
_EXTRACTION_SEMAPHORE = asyncio.Semaphore(CONTENT_EXTRACTION_MAX_CONCURRENCY)


class ContentExtractionError(Exception):
    """Extraction failure that carries a message safe to show the uploader.

    The exception's own text keeps the full technical detail (char counts, timeouts) for
    the logs; `user_message` is the clean, non-technical string the UI shows instead of a
    raw stack-trace-flavoured message (surfaced in routers/retrieval.py process_file)."""

    def __init__(self, message: str, user_message: str):
        super().__init__(message)
        self.user_message = user_message


# Extensions treated as images for the vision-LLM path (see Loader.load).
IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'webp', 'gif', 'bmp', 'tiff', 'tif'}

# Non-PDF types eligible for the Sunway office fast-path (see Loader.load): modern OOXML
# office files + csv. Legacy .doc/.xls/.ppt are excluded -- they already use local loaders
# (and Docling can't parse legacy binary formats anyway).
OFFICE_EXTS = {'docx', 'xlsx', 'pptx', 'csv'}

# Office embedded-image OCR (Sunway). docx/pptx are read text-layer-only by BOTH the fast
# path (unstructured/markitdown) AND Docling's office pipeline (do_ocr=False), so text
# baked into pasted screenshots/figures is missed. We lift embedded images out of the
# OOXML zip and OCR them (via Docling), appending the result. Deliberately CONSTANTS, not
# env vars -- these are finalized knobs; edit + restart to change:
#   _OFFICE_IMAGE_OCR_ENABLED -> master on/off
#   _OFFICE_IMAGE_MIN_PIXELS  -> skip images smaller than this (logos/icons/bullets)
#   _OFFICE_IMAGE_MAX_COUNT   -> cap OCR passes per document (bounds latency / OCR load)
_OFFICE_IMAGE_OCR_ENABLED = True
_OFFICE_IMAGE_MIN_PIXELS = 90_000  # ~300x300 px
_OFFICE_IMAGE_MAX_COUNT = 20
_OFFICE_IMAGE_EXTS = {'png', 'jpg', 'jpeg', 'webp', 'gif', 'bmp', 'tif', 'tiff'}

known_source_ext = [
    'go',
    'py',
    'java',
    'sh',
    'bat',
    'ps1',
    'cmd',
    'js',
    'ts',
    'css',
    'cpp',
    'hpp',
    'h',
    'c',
    'cs',
    'sql',
    'log',
    'ini',
    'pl',
    'pm',
    'r',
    'dart',
    'dockerfile',
    'env',
    'php',
    'hs',
    'hsc',
    'lua',
    'nginxconf',
    'conf',
    'm',
    'mm',
    'plsql',
    'perl',
    'rb',
    'rs',
    'db2',
    'scala',
    'bash',
    'swift',
    'vue',
    'svelte',
    'ex',
    'exs',
    'erl',
    'tsx',
    'jsx',
    'hs',
    'lhs',
    'json',
    'yaml',
    'yml',
    'toml',
]


class ExcelLoader:
    """Fallback Excel loader using pandas when unstructured is not installed."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        import pandas as pd

        text_parts = []
        xls = pd.ExcelFile(self.file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text_parts.append(f'Sheet: {sheet_name}\n{df.to_string(index=False)}')
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class PptxLoader:
    """Fallback PowerPoint loader using python-pptx when unstructured is not installed."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                text_parts.append(f'Slide {i}:\n' + '\n'.join(slide_texts))
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class TikaLoader:
    def __init__(self, url, file_path, mime_type=None, extract_images=None):
        self.url = url
        self.file_path = file_path
        self.mime_type = mime_type

        self.extract_images = extract_images

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            data = f.read()

        if self.mime_type is not None:
            headers = {'Content-Type': self.mime_type}
        else:
            headers = {}

        if self.extract_images == True:
            headers['X-Tika-PDFextractInlineImages'] = 'true'

        endpoint = self.url
        if not endpoint.endswith('/'):
            endpoint += '/'
        endpoint += 'tika/text'

        r = requests.put(
            endpoint,
            data=data,
            headers=headers,
            verify=REQUESTS_VERIFY,
            timeout=(10, CONTENT_EXTRACTION_REQUEST_TIMEOUT),
        )

        if r.ok:
            raw_metadata = r.json()
            text = raw_metadata.get('X-TIKA:content', '<No text content found>').strip()

            if 'Content-Type' in raw_metadata:
                headers['Content-Type'] = raw_metadata['Content-Type']

            log.debug('Tika extracted text: %s', text)

            return [Document(page_content=text, metadata=headers)]
        else:
            raise Exception(f'Error calling Tika: {r.reason}')


class DoclingLoader:
    def __init__(self, url, api_key=None, file_path=None, mime_type=None, params=None):
        self.url = url.rstrip('/')
        self.api_key = api_key
        self.file_path = file_path
        self.mime_type = mime_type

        self.params = params or {}

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            headers = {}
            if self.api_key:
                headers['X-Api-Key'] = f'{self.api_key}'

            r = requests.post(
                f'{self.url}/v1/convert/file',
                files={
                    'files': (
                        self.file_path,
                        f,
                        self.mime_type or 'application/octet-stream',
                    )
                },
                data={
                    'image_export_mode': 'placeholder',
                    **self.params,
                },
                headers=headers,
                verify=AIOHTTP_CLIENT_SESSION_SSL,
                timeout=(10, CONTENT_EXTRACTION_REQUEST_TIMEOUT),
            )
        if r.ok:
            result = r.json()
            document_data = result.get('document', {})
            text = document_data.get('md_content', '<No text content found>')

            metadata = {'Content-Type': self.mime_type} if self.mime_type else {}

            log.debug('Docling extracted text: %s', text)
            return [Document(page_content=text, metadata=metadata)]
        else:
            error_msg = f'Error calling Docling API: {r.reason}'
            if r.text:
                try:
                    error_data = r.json()
                    if 'detail' in error_data:
                        error_msg += f' - {error_data["detail"]}'
                except Exception:
                    error_msg += f' - {r.text}'
            raise Exception(f'Error calling Docling: {error_msg}')


class Loader:
    def __init__(self, engine: str = '', **kwargs):
        self.engine = engine
        self.user = kwargs.get('user', None)
        self.kwargs = kwargs

    def load(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        file_ext = filename.split('.')[-1].lower() if '.' in filename else ''

        # Set when a digital-looking PDF is rerouted to Docling because it embeds images
        # whose baked-in text the fast path would miss; forces OCR on the Docling call.
        force_ocr = False

        # Sunway A/B: fast-path toggles are read at REQUEST time (runtime-mutable via the
        # Admin UI), not from the boot-time env constants, so switching needs no restart.
        ab = get_state()

        # PDF fast-path: born-digital PDFs are extracted by a lightweight loader in
        # milliseconds (pypdf, or MarkItDown when selected), whereas Docling runs
        # layout+OCR models on every page and can take minutes / time out (e.g. a 14MB
        # annual report). Falls back to Docling when the PDF looks scanned (no text layer).
        if ab['pdf_fast_path'] and file_ext == 'pdf' and self.engine == 'docling':
            # Both fast engines read the text layer only (no OCR). Extract with whichever
            # is selected; a non-empty, digital-looking result becomes the fast candidate.
            fast_docs = None
            engine_tag = None
            if ab.get('pdf_engine') == 'markitdown':
                try:
                    from open_webui.retrieval.loaders.markitdown_loader import MarkItDownLoader

                    md_docs = MarkItDownLoader(file_path=file_path, mime_type=file_content_type).load()
                    if md_docs and sum(len((d.page_content or '').strip()) for d in md_docs) >= (
                        RAG_PDF_FAST_PATH_MIN_CHARS_PER_PAGE
                    ):
                        fast_docs = md_docs
                        engine_tag = 'markitdown'
                except Exception as e:
                    log.warning(f'PDF fast-path: markitdown failed ({e}); falling back to {self.engine}')
            else:
                try:
                    pdf_docs = PyPDFLoader(
                        file_path,
                        extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                        mode=self.kwargs.get('PDF_LOADER_MODE', 'page'),
                    ).load()
                    if pdf_docs and self._pdf_looks_digital(pdf_docs):
                        fast_docs = pdf_docs
                        engine_tag = 'pypdf'
                except Exception as e:
                    log.warning(f'PDF fast-path: pypdf failed ({e}); falling back to {self.engine}')

            # Shared image-aware routing (applies to whichever fast engine ran): a digital
            # PDF that ALSO embeds substantial images has baked-in text (screenshots/scanned
            # figures) the text-layer read would miss — reroute to Docling with forced OCR
            # rather than take the fast path.
            if fast_docs:
                if RAG_PDF_IMAGE_ROUTE_ENABLED and self._pdf_has_significant_images(file_path):
                    log.info(
                        f'PDF fast-path: {filename} is digital but image-bearing; routing to '
                        f'{self.engine} with forced OCR to read embedded-image text.'
                    )
                    force_ocr = True
                else:
                    log.info(
                        f'PDF fast-path: born-digital ({len(fast_docs)} pages) via {engine_tag}, skipped {self.engine}'
                    )
                    return [
                        Document(
                            page_content=ftfy.fix_text(doc.page_content),
                            metadata={**doc.metadata, 'processing_engine': engine_tag},
                        )
                        for doc in fast_docs
                    ]
            if not force_ocr:
                log.info(f'PDF fast-path: scanned/low-text — falling back to {self.engine} for OCR')

        # Office fast-path (Sunway A/B): born-digital OOXML office files have a text layer,
        # so a lightweight loader (Unstructured or MarkItDown, chosen at runtime) extracts
        # them far faster than Docling. Falls back to Docling on failure/empty text.
        if ab['office_fast_path'] and file_ext in OFFICE_EXTS and self.engine == 'docling':
            office_docs = self._load_office_fast(filename, file_content_type, file_path, ab['office_engine'])
            if office_docs is not None:
                # docx/pptx are read text-layer only — OCR embedded screenshots/figures
                # and append them (neither the fast path nor Docling's office pipeline
                # does). xlsx/csv rarely embed text-bearing images, so they're skipped.
                if _OFFICE_IMAGE_OCR_ENABLED and file_ext in ('docx', 'pptx'):
                    office_docs = self._append_office_image_ocr(filename, file_path, office_docs)
                return office_docs
            log.info(f'Office fast-path: {filename} empty/failed — falling back to {self.engine}')

        # Image vision path (Sunway): route images to a vision LLM so text-only chat
        # models get non-text visual understanding OCR can't provide (see env
        # RAG_IMAGE_VISION_LLM_*). Optionally combined with Docling OCR for faithful
        # text. Falls through to the normal engine when the vision LLM is unconfigured.
        if file_ext in IMAGE_EXTS:
            vision_ready = bool(RAG_IMAGE_VISION_LLM_BASE_URL and RAG_IMAGE_VISION_LLM_MODEL)
            log.info(
                f'Image {filename}: vision LLM configured={vision_ready} '
                f'(base_url_set={bool(RAG_IMAGE_VISION_LLM_BASE_URL)}, '
                f'model={RAG_IMAGE_VISION_LLM_MODEL or "<unset>"}, engine={self.engine!r})'
            )
            if vision_ready:
                return self._load_image_with_vision(filename, file_content_type, file_path)

        # Image-aware OCR, decoupled from the fast-path so it ALSO applies when a file
        # reaches Docling directly (fast-path off/failed) -- otherwise embedded-image text
        # is silently missed in single-engine mode. Shared across all extraction routes.
        # PDF: force OCR when the PDF embeds substantial images, so text-layer pages with
        # pasted screenshots still get OCR'd (the fast-path branch sets force_ocr itself;
        # `not force_ocr` avoids re-checking a PDF it already decided).
        if (
            file_ext == 'pdf'
            and self.engine == 'docling'
            and not force_ocr
            and RAG_PDF_IMAGE_ROUTE_ENABLED
            and self._pdf_has_significant_images(file_path)
        ):
            log.info(f'PDF: {filename} is image-bearing (outside fast-path); forcing OCR for embedded-image text.')
            force_ocr = True

        loader = self._get_loader(filename, file_content_type, file_path, force_ocr=force_ocr)
        docs = loader.load()
        result = [Document(page_content=ftfy.fix_text(doc.page_content), metadata=doc.metadata) for doc in docs]

        # Office: OCR embedded images for docx/pptx that reached Docling directly. The
        # office fast-path branch already handles (and returns early for) the fast-path
        # success case, so this covers only the fast-path off/failed route -- no double.
        if _OFFICE_IMAGE_OCR_ENABLED and self.engine == 'docling' and file_ext in ('docx', 'pptx'):
            result = self._append_office_image_ocr(filename, file_path, result)

        return result

    def _load_office_fast(
        self, filename: str, file_content_type: str, file_path: str, engine: str
    ) -> list[Document] | None:
        """Extract a born-digital OOXML office file via a lightweight loader.

        engine='unstructured' uses the bundled langchain Unstructured loaders;
        'markitdown' uses Microsoft MarkItDown. Returns tagged Documents, or None when the
        loader errors or yields no text (the caller then falls back to Docling).
        """
        file_ext = filename.split('.')[-1].lower()
        try:
            if engine == 'markitdown':
                from open_webui.retrieval.loaders.markitdown_loader import MarkItDownLoader

                docs = MarkItDownLoader(file_path=file_path, mime_type=file_content_type).load()
                tag = 'markitdown'
            elif file_ext == 'docx':
                from langchain_community.document_loaders import UnstructuredWordDocumentLoader

                docs = UnstructuredWordDocumentLoader(file_path).load()
                tag = 'unstructured'
            elif file_ext == 'xlsx':
                from langchain_community.document_loaders import UnstructuredExcelLoader

                docs = UnstructuredExcelLoader(file_path).load()
                tag = 'unstructured'
            elif file_ext == 'pptx':
                from langchain_community.document_loaders import UnstructuredPowerPointLoader

                docs = UnstructuredPowerPointLoader(file_path).load()
                tag = 'unstructured'
            elif file_ext == 'csv':
                docs = CSVLoader(file_path, encoding=self._detect_text_encoding(file_path)).load()
                tag = 'csvloader'
            else:
                return None
        except Exception as e:
            log.warning(f'Office fast-path ({engine}) failed for {filename}: {e}; falling back')
            return None

        total_chars = sum(len((doc.page_content or '').strip()) for doc in docs)
        if not docs or total_chars == 0:
            return None

        log.info(f'Office fast-path: {filename} via {tag} ({total_chars} chars), skipped {self.engine}')
        return [
            Document(
                page_content=ftfy.fix_text(doc.page_content),
                metadata={**doc.metadata, 'processing_engine': tag},
            )
            for doc in docs
        ]

    def _append_office_image_ocr(self, filename: str, file_path: str, office_docs: list[Document]) -> list[Document]:
        """Append OCR of a docx/pptx's embedded images (screenshots/figures) to its
        text-layer extraction. Best-effort: any failure leaves office_docs unchanged."""
        import tempfile

        try:
            with tempfile.TemporaryDirectory() as tmp:
                image_paths = self._extract_office_images(file_path, tmp)
                if not image_paths:
                    return office_docs
                ocr_text = self._ocr_office_images(image_paths)
        except Exception as e:
            log.warning(f'Office image OCR failed for {filename}: {e}')
            return office_docs

        if not ocr_text:
            return office_docs

        log.info(
            f'Office image OCR: {filename} — appended {len(ocr_text)} chars from {len(image_paths)} embedded image(s)'
        )
        return office_docs + [
            Document(
                page_content=ftfy.fix_text(ocr_text),
                metadata={'processing_engine': 'office-img-ocr'},
            )
        ]

    def _extract_office_images(self, file_path: str, dest_dir: str) -> list[str]:
        """Lift substantial embedded raster images out of an OOXML (docx/pptx) zip.

        OOXML files are zips; images live under word/media/ or ppt/media/. Keeps only
        images >= _OFFICE_IMAGE_MIN_PIXELS (skips logos/icons/bullets), up to
        _OFFICE_IMAGE_MAX_COUNT (bounds OCR cost). Cheap dimension read via Pillow; never
        raises -- returns the paths written into dest_dir (possibly empty)."""
        import os
        import zipfile
        from io import BytesIO

        from PIL import Image

        paths: list[str] = []
        try:
            with zipfile.ZipFile(file_path) as z:
                for name in z.namelist():
                    if len(paths) >= _OFFICE_IMAGE_MAX_COUNT:
                        break
                    if '/media/' not in name.lower():
                        continue
                    ext = name.rsplit('.', 1)[-1].lower() if '.' in name else ''
                    if ext not in _OFFICE_IMAGE_EXTS:
                        continue
                    data = z.read(name)
                    try:
                        with Image.open(BytesIO(data)) as im:
                            width, height = im.size
                    except Exception:
                        continue
                    if width * height < _OFFICE_IMAGE_MIN_PIXELS:
                        continue
                    out_path = os.path.join(dest_dir, f'office_img_{len(paths)}.{ext}')
                    with open(out_path, 'wb') as f:
                        f.write(data)
                    paths.append(out_path)
        except Exception as e:
            log.warning(f'Office image extraction failed for {file_path}: {e}')
        return paths

    def _ocr_office_images(self, image_paths: list[str]) -> str:
        """OCR each extracted office image via Docling and return the combined text.
        Sequential (bounded by _OFFICE_IMAGE_MAX_COUNT); per-image failures are skipped."""
        import mimetypes

        server_url = self.kwargs.get('DOCLING_SERVER_URL')
        if not server_url:
            log.warning('Office image OCR: DOCLING_SERVER_URL not set; skipping')
            return ''

        params = self.kwargs.get('DOCLING_PARAMS', {})
        if not isinstance(params, dict):
            try:
                params = json.loads(params)
            except json.JSONDecodeError:
                params = {}
        # These images ARE the content to read -- ensure OCR is on.
        params = {**params, 'do_ocr': True}

        sections: list[str] = []
        for path in image_paths:
            mime = mimetypes.guess_type(path)[0] or 'image/png'
            try:
                docs = DoclingLoader(
                    url=server_url,
                    api_key=self.kwargs.get('DOCLING_API_KEY', None),
                    file_path=path,
                    mime_type=mime,
                    params=params,
                ).load()
                text = '\n'.join((d.page_content or '').strip() for d in docs).strip()
                if text:
                    sections.append(text)
            except Exception as e:
                log.warning(f'Office image OCR: failed on {path}: {e}')
        return '\n\n'.join(sections)

    def _load_image_with_vision(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        """Extract an image via the vision LLM, optionally prepending Docling OCR.

        Gives a text-only chat model both faithful transcribed text (Docling, when
        COMBINE_OCR is on and the engine is Docling) and a description of non-text
        visuals (the vision LLM) that OCR alone cannot produce. Returns a single
        combined Document.
        """
        from open_webui.retrieval.loaders.vision_llm import (
            DEFAULT_VISION_DESCRIBE_PROMPT,
            DEFAULT_VISION_PROMPT,
            VisionLLMLoader,
        )

        extra_body = {}
        if RAG_IMAGE_VISION_LLM_EXTRA_BODY:
            try:
                extra_body = json.loads(RAG_IMAGE_VISION_LLM_EXTRA_BODY)
            except json.JSONDecodeError:
                log.error('Invalid RAG_IMAGE_VISION_LLM_EXTRA_BODY (expected a JSON object); ignoring')

        sections = []

        # Faithful OCR first (best for text-dense images; a VLM can misread exact text).
        # Best-effort: on failure, fall back to the vision LLM alone.
        ocr_text = ''
        combine_ocr = (
            RAG_IMAGE_VISION_LLM_COMBINE_OCR and self.engine == 'docling' and self.kwargs.get('DOCLING_SERVER_URL')
        )
        if combine_ocr:
            try:
                params = self.kwargs.get('DOCLING_PARAMS', {})
                if not isinstance(params, dict):
                    try:
                        params = json.loads(params)
                    except json.JSONDecodeError:
                        params = {}
                docling_docs = DoclingLoader(
                    url=self.kwargs.get('DOCLING_SERVER_URL'),
                    api_key=self.kwargs.get('DOCLING_API_KEY', None),
                    file_path=file_path,
                    mime_type=file_content_type,
                    params=params,
                ).load()
                ocr_text = '\n'.join(doc.page_content for doc in docling_docs).strip()
            except Exception as e:
                log.warning(f'Vision image path: Docling OCR step failed ({e}); using vision LLM only')

        if ocr_text:
            sections.append(f'## Extracted text\n\n{ocr_text}')

        # Vision LLM prompt: describe-only whenever the OCR step RAN -- OCR is authoritative
        # for text even when it found NONE. Keying this off `ocr_text` instead would flip a
        # no-text image (e.g. a portrait photo) back to the transcribe-everything prompt,
        # which is exactly when a small VLM invents text ("paz" on a photo of a face) --
        # garbage that then pollutes the vector DB, wastes context, and can drag the chat
        # reply into the hallucinated token's language. Full extract only when no OCR step
        # ran at all (the vision LLM is then the sole reader and must transcribe).
        if RAG_IMAGE_VISION_LLM_PROMPT:
            prompt = RAG_IMAGE_VISION_LLM_PROMPT
        else:
            prompt = DEFAULT_VISION_DESCRIBE_PROMPT if combine_ocr else DEFAULT_VISION_PROMPT

        vision_docs = VisionLLMLoader(
            base_url=RAG_IMAGE_VISION_LLM_BASE_URL,
            model=RAG_IMAGE_VISION_LLM_MODEL,
            file_path=file_path,
            api_key=RAG_IMAGE_VISION_LLM_API_KEY,
            mime_type=file_content_type,
            prompt=prompt,
            max_tokens=RAG_IMAGE_VISION_LLM_MAX_TOKENS,
            extra_body=extra_body,
        ).load()
        vision_text = '\n'.join(doc.page_content for doc in vision_docs).strip()

        # Drop the describe-only sentinel so pure-text images don't get an empty section.
        if vision_text.lower().strip('.() ') == 'no additional visual content':
            vision_text = ''

        if vision_text:
            sections.append(f'## Visual description\n\n{vision_text}' if ocr_text else vision_text)

        combined = '\n\n'.join(sections).strip()
        if not combined:
            combined = 'No content could be extracted from the image.'

        return [
            Document(
                page_content=ftfy.fix_text(combined),
                metadata={'Content-Type': file_content_type, 'processing_engine': 'vision-llm'},
            )
        ]

    def _pdf_looks_digital(self, docs: list[Document]) -> bool:
        """True if pypdf extracted enough text per page to treat the PDF as
        born-digital; scanned/image-only PDFs yield ~0 chars/page. Threshold is
        RAG_PDF_FAST_PATH_MIN_CHARS_PER_PAGE."""
        if not docs:
            return False
        total_chars = sum(len((doc.page_content or '').strip()) for doc in docs)
        return (total_chars / len(docs)) >= RAG_PDF_FAST_PATH_MIN_CHARS_PER_PAGE

    def _pdf_has_significant_images(self, file_path: str) -> bool:
        """True if the PDF embeds enough large raster images that its text layer alone
        likely misses baked-in text (pasted screenshots, scanned figures) — the signal to
        reroute a 'digital' PDF to Docling with forced OCR.

        Cheap: reads each image XObject's /Width and /Height from the page resources
        WITHOUT decoding pixels, and short-circuits once the count threshold is met. Tuned
        by RAG_PDF_IMAGE_MIN_PIXELS (skip small logos/icons) and RAG_PDF_IMAGE_MIN_COUNT.
        Detection can only see image SIZE, not whether an image contains text, so it errs
        toward rerouting. Never raises — on any error it declines to reroute (fast path)."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            substantial = 0
            for page in reader.pages:
                resources = page.get('/Resources')
                if not resources:
                    continue
                xobjects = resources.get('/XObject')
                if not xobjects:
                    continue
                try:
                    xobjects = xobjects.get_object()
                except Exception:
                    continue
                for name in xobjects:
                    try:
                        xobj = xobjects[name].get_object()
                    except Exception:
                        continue
                    if xobj.get('/Subtype') != '/Image':
                        continue
                    width = int(xobj.get('/Width', 0) or 0)
                    height = int(xobj.get('/Height', 0) or 0)
                    if width * height >= RAG_PDF_IMAGE_MIN_PIXELS:
                        substantial += 1
                        if substantial >= RAG_PDF_IMAGE_MIN_COUNT:
                            return True
            return False
        except Exception as e:
            log.warning(f'PDF image-detection failed for {file_path}: {e}; skipping image reroute')
            return False

    async def aload(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        """
        Async wrapper around `load`.

        Document loaders dispatched by `_get_loader` (PyMuPDF, Unstructured,
        python-docx, Tika, etc.) are uniformly synchronous and CPU/IO-bound.
        Calling `load` directly from an async handler would block the event
        loop for the entire parse — minutes for large PDFs. This offloads
        the work to a worker thread so the loop stays responsive.
        """
        async with _EXTRACTION_SEMAPHORE:
            # (A) Wall-clock backstop. The in-process loaders have no timeout of their
            # own, so a pathological file could pin an extraction slot indefinitely.
            # On timeout we release the slot and surface a clean error; the orphaned
            # worker thread can't be force-killed but no longer blocks new uploads.
            try:
                docs = await asyncio.wait_for(
                    asyncio.to_thread(self.load, filename, file_content_type, file_path),
                    timeout=CONTENT_EXTRACTION_TIMEOUT,
                )
            except asyncio.TimeoutError:
                raise ContentExtractionError(
                    f'Content extraction for {filename} exceeded the '
                    f'{CONTENT_EXTRACTION_TIMEOUT}s limit and was aborted.',
                    'This file took too long to process and was stopped. Try a smaller or simpler file.',
                )

            # (B) Decompression-bomb guard. A small archive-based file (docx/xlsx/pptx
            # are ZIP containers) or crafted PDF can expand to gigabytes of text; reject
            # before it reaches chunking/embedding/the vector DB.
            if CONTENT_EXTRACTION_MAX_OUTPUT_CHARS:
                total_chars = sum(len(doc.page_content or '') for doc in docs)
                if total_chars > CONTENT_EXTRACTION_MAX_OUTPUT_CHARS:
                    raise ContentExtractionError(
                        f'Extracted content from {filename} ({total_chars} chars) exceeds the '
                        f'{CONTENT_EXTRACTION_MAX_OUTPUT_CHARS}-char limit and was rejected.',
                        'This file is too large to process. Try a smaller or compressed file, '
                        'or upload only the content you need.',
                    )

            return docs

    def _is_text_file(self, file_ext: str, file_content_type: str) -> bool:
        return file_ext in known_source_ext or (
            file_content_type
            and file_content_type.find('text/') >= 0
            # Avoid text/html files being detected as text
            and not file_content_type.find('html') >= 0
        )

    def _detect_text_encoding(self, file_path: str) -> str:
        """Detect the encoding of a text file with CJK-aware fallbacks.

        Langchain's ``TextLoader`` uses chardet internally when
        ``autodetect_encoding=True``, but chardet frequently misidentifies
        CJK encodings (e.g. GB18030 detected as GB2312 or even Cyrillic).
        This method replaces that by:

        1. Trying UTF-8 first (fast path for the vast majority of files).
        2. Using chardet as a *hint* to prioritise the right CJK codec
           family, but mapping subset names to their superset
           (e.g. GB2312 → gb18030).
        3. Validating that decoded text actually contains CJK characters,
           guarding against codecs that "succeed" but produce garbage.
        4. Falling back to latin-1 (always valid, ftfy fixes mojibake later).
        """
        try:
            with open(file_path, 'rb') as f:
                raw = f.read()
        except OSError:
            return 'utf-8'

        if not raw:
            return 'utf-8'

        # Fast path: most files are UTF-8
        try:
            raw.decode('utf-8')
            return 'utf-8'
        except UnicodeDecodeError:
            pass

        # Use chardet as a hint, not as ground truth
        import chardet

        detected = chardet.detect(raw)
        detected_enc = (detected.get('encoding') or '').lower().replace('-', '').replace('_', '')

        # Map chardet's detected encoding to the correct superset codec.
        # chardet often reports GB2312 for content that is actually GB18030;
        # GB18030 is a strict superset of both GB2312 and GBK.
        _ENC_FAMILY = {
            'gb2312': 'gb18030',
            'gb18030': 'gb18030',
            'gbk': 'gb18030',
            'big5': 'big5',
            'euckr': 'euc-kr',
            'eucjp': 'euc-jp',
            'iso2022jp': 'euc-jp',
            'shiftjis': 'shift_jis',
        }

        # Build priority list: chardet-hinted codec first, then remaining CJK
        base_order = ['gb18030', 'big5', 'euc-kr', 'euc-jp']
        hinted = _ENC_FAMILY.get(detected_enc)
        if hinted and hinted in base_order:
            ordered = [hinted] + [e for e in base_order if e != hinted]
        else:
            ordered = base_order

        for enc in ordered:
            try:
                text = raw.decode(enc)
                if text.strip() and self._has_cjk_characters(text):
                    log.info(
                        'Detected encoding %s for %s (chardet guessed %s)',
                        enc,
                        file_path,
                        detected.get('encoding'),
                    )
                    return enc
            except (UnicodeDecodeError, LookupError):
                continue

        # If chardet gave a non-CJK answer that isn't in our family map,
        # try it directly — it might be a valid Western encoding.
        chardet_encoding = detected.get('encoding')
        if chardet_encoding:
            try:
                raw.decode(chardet_encoding)
                log.info(
                    'Using chardet-detected encoding %s for %s',
                    chardet_encoding,
                    file_path,
                )
                return chardet_encoding
            except (UnicodeDecodeError, LookupError):
                pass

        # latin-1 is the ultimate fallback: every byte 0x00–0xFF is valid.
        # ftfy.fix_text() (applied downstream) repairs most mojibake that
        # results from treating Windows-1252 content as Latin-1.
        log.info('Falling back to latin-1 encoding for %s', file_path)
        return 'latin-1'

    @staticmethod
    def _has_cjk_characters(text: str, threshold: float = 0.05) -> bool:
        """Check if decoded text contains a meaningful proportion of CJK characters.

        This guards against codecs that technically "succeed" but decode the
        bytes into wrong Unicode codepoints (e.g. PUA chars, random symbols).
        A genuine CJK document should have at least ``threshold`` fraction of
        its non-whitespace characters in CJK Unicode blocks.
        """
        if not text:
            return False

        cjk_count = 0
        total = 0
        for ch in text:
            if ch.isspace():
                continue
            total += 1
            cp = ord(ch)
            if (
                0x4E00 <= cp <= 0x9FFF  # CJK Unified Ideographs
                or 0x3400 <= cp <= 0x4DBF  # CJK Extension A
                or 0x20000 <= cp <= 0x2A6DF  # CJK Extension B
                or 0x2A700 <= cp <= 0x2B73F  # CJK Extension C
                or 0x2B740 <= cp <= 0x2B81F  # CJK Extension D
                or 0xF900 <= cp <= 0xFAFF  # CJK Compatibility Ideographs
                or 0x3000 <= cp <= 0x303F  # CJK Symbols and Punctuation
                or 0x3040 <= cp <= 0x309F  # Hiragana
                or 0x30A0 <= cp <= 0x30FF  # Katakana
                or 0xAC00 <= cp <= 0xD7AF  # Hangul Syllables
                or 0xFF00 <= cp <= 0xFFEF  # Halfwidth and Fullwidth Forms
            ):
                cjk_count += 1

        if total == 0:
            return False

        return (cjk_count / total) >= threshold

    def _get_loader(self, filename: str, file_content_type: str, file_path: str, force_ocr: bool = False):
        file_ext = filename.split('.')[-1].lower()

        if (
            self.engine == 'external'
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL')
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY')
        ):
            loader = ExternalDocumentLoader(
                file_path=file_path,
                url=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL'),
                api_key=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY'),
                mime_type=file_content_type,
                user=self.user,
            )
        elif self.engine == 'tika' and self.kwargs.get('TIKA_SERVER_URL'):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                loader = TikaLoader(
                    url=self.kwargs.get('TIKA_SERVER_URL'),
                    file_path=file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                )
        elif (
            self.engine == 'datalab_marker'
            and self.kwargs.get('DATALAB_MARKER_API_KEY')
            and file_ext
            in [
                'pdf',
                'xls',
                'xlsx',
                'ods',
                'doc',
                'docx',
                'odt',
                'ppt',
                'pptx',
                'odp',
                'html',
                'epub',
                'png',
                'jpeg',
                'jpg',
                'webp',
                'gif',
                'tiff',
            ]
        ):
            api_base_url = self.kwargs.get('DATALAB_MARKER_API_BASE_URL', '')
            if not api_base_url or api_base_url.strip() == '':
                api_base_url = 'https://www.datalab.to/api/v1/marker'  # https://github.com/open-webui/open-webui/pull/16867#issuecomment-3218424349

            loader = DatalabMarkerLoader(
                file_path=file_path,
                api_key=self.kwargs['DATALAB_MARKER_API_KEY'],
                api_base_url=api_base_url,
                additional_config=self.kwargs.get('DATALAB_MARKER_ADDITIONAL_CONFIG'),
                use_llm=self.kwargs.get('DATALAB_MARKER_USE_LLM', False),
                skip_cache=self.kwargs.get('DATALAB_MARKER_SKIP_CACHE', False),
                force_ocr=self.kwargs.get('DATALAB_MARKER_FORCE_OCR', False),
                paginate=self.kwargs.get('DATALAB_MARKER_PAGINATE', False),
                strip_existing_ocr=self.kwargs.get('DATALAB_MARKER_STRIP_EXISTING_OCR', False),
                disable_image_extraction=self.kwargs.get('DATALAB_MARKER_DISABLE_IMAGE_EXTRACTION', False),
                format_lines=self.kwargs.get('DATALAB_MARKER_FORMAT_LINES', False),
                output_format=self.kwargs.get('DATALAB_MARKER_OUTPUT_FORMAT', 'markdown'),
            )
        elif (
            self.engine == 'docling'
            and self.kwargs.get('DOCLING_SERVER_URL')
            # Docling only parses 2007+ (OOXML) Office files, so legacy binary formats
            # must fall through to the local loaders below: .xls works via
            # unstructured -> pandas/xlrd; .doc/.ppt additionally require LibreOffice
            # (soffice) on PATH, which the container image does not ship. Match on
            # extension, not content type: Windows browsers with Excel installed
            # report .csv as application/vnd.ms-excel, and csv must stay on Docling.
            and file_ext not in ('doc', 'xls', 'ppt')
        ):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                # Build params for DoclingLoader
                params = self.kwargs.get('DOCLING_PARAMS', {})
                if not isinstance(params, dict):
                    try:
                        params = json.loads(params)
                    except json.JSONDecodeError:
                        log.error('Invalid DOCLING_PARAMS format, expected JSON object')
                        params = {}
                if force_ocr:
                    # This PDF was rerouted here for its embedded images; OCR every page so
                    # their baked-in text is read (the selective default skips pages that
                    # already have a text layer, missing text inside pasted screenshots).
                    params = {**params, 'force_ocr': True}

                loader = DoclingLoader(
                    url=self.kwargs.get('DOCLING_SERVER_URL'),
                    api_key=self.kwargs.get('DOCLING_API_KEY', None),
                    file_path=file_path,
                    mime_type=file_content_type,
                    params=params,
                )
        elif (
            self.engine == 'document_intelligence'
            and self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT') != ''
            and (
                file_ext in ['pdf', 'docx', 'ppt', 'pptx']
                or file_content_type
                in [
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'application/vnd.ms-powerpoint',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                ]
            )
        ):
            if self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY') != '':
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    api_key=self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY'),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
            else:
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    azure_credential=DefaultAzureCredential(),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
        elif self.engine == 'mineru' and file_ext in self.kwargs.get('MINERU_FILE_EXTENSIONS', ['pdf']):
            mineru_timeout = self.kwargs.get('MINERU_API_TIMEOUT', 300)
            if mineru_timeout:
                try:
                    mineru_timeout = int(mineru_timeout)
                except ValueError:
                    mineru_timeout = 300

            loader = MinerULoader(
                file_path=file_path,
                api_mode=self.kwargs.get('MINERU_API_MODE', 'local'),
                api_url=self.kwargs.get('MINERU_API_URL', 'http://localhost:8000'),
                api_key=self.kwargs.get('MINERU_API_KEY', ''),
                params=self.kwargs.get('MINERU_PARAMS', {}),
                timeout=mineru_timeout,
            )
        elif (
            self.engine == 'mistral_ocr'
            and self.kwargs.get('MISTRAL_OCR_API_KEY') != ''
            and file_ext in ['pdf']  # Mistral OCR currently only supports PDF and images
        ):
            loader = MistralLoader(
                base_url=self.kwargs.get('MISTRAL_OCR_API_BASE_URL'),
                api_key=self.kwargs.get('MISTRAL_OCR_API_KEY'),
                file_path=file_path,
            )
        elif self.engine == 'paddleocr_vl' and self.kwargs.get('PADDLEOCR_VL_TOKEN') != '':
            loader = PaddleOCRVLLoader(
                api_url=self.kwargs.get('PADDLEOCR_VL_BASE_URL'),
                token=self.kwargs.get('PADDLEOCR_VL_TOKEN'),
                file_path=file_path,
            )
        else:
            if file_ext == 'pdf':
                loader = PyPDFLoader(
                    file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                    mode=self.kwargs.get('PDF_LOADER_MODE', 'page'),
                )
            elif file_ext == 'csv':
                loader = CSVLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext == 'rst':
                try:
                    from langchain_community.document_loaders import UnstructuredRSTLoader

                    loader = UnstructuredRSTLoader(file_path, mode='elements')
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .rst file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext == 'xml':
                try:
                    from langchain_community.document_loaders import UnstructuredXMLLoader

                    loader = UnstructuredXMLLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .xml file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext in ['htm', 'html']:
                loader = BSHTMLLoader(file_path, open_encoding='unicode_escape')
            elif file_ext == 'md':
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_content_type == 'application/epub+zip':
                try:
                    from langchain_community.document_loaders import UnstructuredEPubLoader

                    loader = UnstructuredEPubLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .epub files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif (
                file_content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                or file_ext == 'docx'
            ):
                loader = Docx2txtLoader(file_path)
            elif file_ext == 'doc' or file_content_type == 'application/msword':
                try:
                    from langchain_community.document_loaders import UnstructuredWordDocumentLoader

                    loader = UnstructuredWordDocumentLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .doc files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif file_content_type in [
                'application/vnd.ms-excel',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            ] or file_ext in ['xls', 'xlsx']:
                if file_ext == 'xls':
                    # Legacy BIFF .xls goes straight to pandas/xlrd: unstructured's
                    # partition_xlsx pre-scans with msoffcrypto, which crashes
                    # (struct.error) on minimal xlwt-style streams that xlrd reads fine.
                    loader = ExcelLoader(file_path)
                else:
                    try:
                        from langchain_community.document_loaders import UnstructuredExcelLoader

                        loader = UnstructuredExcelLoader(file_path)
                    except ImportError:
                        log.warning(
                            "The 'unstructured' package is not installed. "
                            'Falling back to pandas for Excel file loading. '
                            'Install unstructured for better results: pip install unstructured'
                        )
                        loader = ExcelLoader(file_path)
            elif file_content_type in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            ] or file_ext in ['ppt', 'pptx']:
                try:
                    from langchain_community.document_loaders import UnstructuredPowerPointLoader

                    loader = UnstructuredPowerPointLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to python-pptx for PowerPoint file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = PptxLoader(file_path)
            elif file_ext == 'msg':
                loader = OutlookMessageLoader(file_path)
            elif file_ext == 'odt':
                try:
                    from langchain_community.document_loaders import UnstructuredODTLoader

                    loader = UnstructuredODTLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .odt files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))

        return loader
