"""Extraction stress harness — measures how long the built-in ("lightweight")
loaders take to extract text from large/many common files (csv, xlsx, pdf, pptx).

Generates synthetic files at three scales each, runs them through the real
`Loader` at the chosen engine, and prints size + extraction time + char count.

  python scripts/extraction_stress.py                 # lightweight built-ins (engine='')
  python scripts/extraction_stress.py --engine docling # needs docling-serve up + DOCLING_SERVER_URL

Run inside the app env (WEBUI_SECRET_KEY etc. set). docx is skipped (no
python-docx installed to generate one); bring a real .docx sample to test it.
"""
import argparse
import csv as csvmod
import os
import shutil
import tempfile
import time
import tracemalloc

from open_webui.retrieval.loaders.main import Loader


def gen_csv(path, rows):
    with open(path, 'w', newline='', encoding='utf-8') as f:
        w = csvmod.writer(f)
        w.writerow([f'col{i}' for i in range(10)])
        for r in range(rows):
            w.writerow([f'r{r}c{c}_lorem_ipsum_sample_data' for c in range(10)])


def gen_xlsx(path, rows):
    from openpyxl import Workbook

    wb = Workbook(write_only=True)
    ws = wb.create_sheet()
    ws.append([f'col{i}' for i in range(10)])
    for r in range(rows):
        ws.append([f'r{r}c{c}_data' for c in range(10)])
    wb.save(path)


def gen_pdf(path, pages):
    from fpdf import FPDF

    pdf = FPDF()
    para = 'The quick brown fox jumps over the lazy dog. ' * 20
    for p in range(pages):
        pdf.add_page()
        pdf.set_font('Helvetica', size=10)
        pdf.multi_cell(0, 5, f'Page {p}\n{para}')
    pdf.output(path)


def gen_pptx(path, slides):
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for s in range(slides):
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Pt(50), Pt(50), Pt(500), Pt(300)).text_frame
        tf.text = f'Slide {s}: ' + ('content text ' * 30)
    prs.save(path)


JOBS = [
    ('csv', '.csv', 'text/csv', gen_csv, [(10_000, '10k rows'), (100_000, '100k rows'), (500_000, '500k rows')]),
    ('xlsx', '.xlsx', None, gen_xlsx, [(5_000, '5k rows'), (25_000, '25k rows'), (100_000, '100k rows')]),
    ('pdf', '.pdf', 'application/pdf', gen_pdf, [(50, '50 pages'), (500, '500 pages'), (1500, '1500 pages')]),
    ('pptx', '.pptx', None, gen_pptx, [(50, '50 slides'), (250, '250 slides'), (800, '800 slides')]),
]


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--engine', default='', help="extraction engine ('' = lightweight built-ins; 'docling' etc.)")
    args = ap.parse_args()

    kwargs = {}
    if args.engine == 'docling':
        kwargs['DOCLING_SERVER_URL'] = os.getenv('DOCLING_SERVER_URL', 'http://localhost:5001')
    loader = Loader(engine=args.engine, **kwargs)

    tmp = tempfile.mkdtemp(prefix='extr_stress_')
    print(f"engine={args.engine or '(lightweight built-ins)'}  tmp={tmp}")
    print(f"\n{'type':5} {'scale':11} {'size':>9} {'gen_s':>7} {'extract_s':>10} {'mem_MB':>8} {'chars':>12}")
    print('-' * 70)
    try:
        for typ, ext, ct, gen, scales in JOBS:
            for count, label in scales:
                path = os.path.join(tmp, f'stress_{typ}_{count}{ext}')
                t0 = time.perf_counter()
                try:
                    gen(path, count)
                except Exception as e:
                    print(f'{typ:5} {label:11} GEN-FAIL {type(e).__name__}: {e}')
                    continue
                gen_s = time.perf_counter() - t0
                mb = os.path.getsize(path) / 1024 / 1024
                tracemalloc.start()
                t1 = time.perf_counter()
                try:
                    docs = loader.load(os.path.basename(path), ct, path)
                    ex_s = time.perf_counter() - t1
                    mem_mb = tracemalloc.get_traced_memory()[1] / 1024 / 1024
                    chars = sum(len(d.page_content) for d in docs)
                    print(f'{typ:5} {label:11} {mb:7.1f}MB {gen_s:7.1f} {ex_s:10.2f} {mem_mb:8.1f} {chars:12,}')
                except Exception as e:
                    print(f'{typ:5} {label:11} {mb:7.1f}MB  EXTRACT-FAIL {type(e).__name__}: {str(e)[:60]}')
                finally:
                    tracemalloc.stop()
                    if os.path.exists(path):
                        os.remove(path)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == '__main__':
    main()
